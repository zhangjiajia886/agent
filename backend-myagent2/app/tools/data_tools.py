"""
数据库 & 文档工具：mysql_query, mysql_schema, redis_cmd, sqlite_query,
                   excel_read, excel_write, word_read, word_write,
                   ppt_read, ppt_write, md_to_html
"""
from __future__ import annotations

import json
import os
from typing import Any

from .base import BaseTool


# ─────────────────────────────────────────────────────────────────────────────
# G. 数据库 & 存储
# ─────────────────────────────────────────────────────────────────────────────

def _mysql_defaults(arguments: dict[str, Any]) -> dict[str, Any]:
    """Fill missing MySQL connection params from MYSQL_* env vars."""
    env = os.environ
    return {
        "host":     arguments.get("host")     or env.get("MYSQL_HOST", "127.0.0.1"),
        "port":     int(arguments.get("port")  or env.get("MYSQL_PORT", 3306)),
        "user":     arguments.get("user")     or env.get("MYSQL_USER", "root"),
        "password": arguments.get("password") if arguments.get("password") is not None
                    else env.get("MYSQL_PASSWORD", ""),
        "database": arguments.get("database") or env.get("MYSQL_DATABASE", ""),
    }


class MySQLQueryTool(BaseTool):
    name = "mysql_query"
    description = "连接 MySQL 执行 SQL 语句。连接参数可省略，自动读取 MYSQL_* 环境变量"
    category = "database"
    risk_level = "high"

    @property
    def input_schema(self) -> dict:
        return {
            "type": "object",
            "properties": {
                "host":     {"type": "string", "description": "默认读 MYSQL_HOST 环境变量"},
                "port":     {"type": "integer", "description": "默认读 MYSQL_PORT 环境变量"},
                "user":     {"type": "string",  "description": "默认读 MYSQL_USER 环境变量"},
                "password": {"type": "string",  "description": "默认读 MYSQL_PASSWORD 环境变量"},
                "database": {"type": "string",  "description": "默认读 MYSQL_DATABASE 环境变量"},
                "sql":      {"type": "string",  "description": "要执行的 SQL 语句"},
                "max_rows": {"type": "integer", "default": 200},
            },
            "required": ["sql"],
        }

    async def run(self, arguments: dict[str, Any]) -> dict[str, Any]:
        import aiomysql
        cfg = _mysql_defaults(arguments)
        host, port, user, password, db = (
            cfg["host"], cfg["port"], cfg["user"], cfg["password"], cfg["database"]
        )
        sql = (arguments.get("sql") or arguments.get("query") or arguments.get("statement") or "").strip()
        if not sql:
            return {"error": "缺少 sql 参数"}
        max_rows = int(arguments.get("max_rows", 200))
        # Split on semicolons to support multi-statement batches
        stmts = [s.strip() for s in sql.split(";") if s.strip()]
        try:
            conn = await aiomysql.connect(
                host=host, port=port, user=user, password=password, db=db,
                charset="utf8mb4", autocommit=True,
            )
            results = []
            async with conn.cursor(aiomysql.DictCursor) as cur:
                for stmt in stmts:
                    await cur.execute(stmt)
                    up = stmt.upper().lstrip()
                    if up.startswith("SELECT") or up.startswith("SHOW") or up.startswith("DESC"):
                        rows = await cur.fetchmany(max_rows)
                        results.append({"sql": stmt[:120], "rows": [dict(r) for r in rows], "count": len(rows)})
                    else:
                        results.append({"sql": stmt[:120], "affected_rows": cur.rowcount, "ok": True})
            conn.close()
            return results[0] if len(results) == 1 else {"results": results}
        except Exception as e:
            return {"error": str(e)}


class MySQLSchemaTool(BaseTool):
    name = "mysql_schema"
    description = "查看 MySQL 数据库的表列表或某张表字段结构。连接参数可省略，自动读取 MYSQL_* 环境变量"
    category = "database"
    risk_level = "low"

    @property
    def input_schema(self) -> dict:
        return {
            "type": "object",
            "properties": {
                "host":     {"type": "string", "description": "默认读 MYSQL_HOST 环境变量"},
                "port":     {"type": "integer", "description": "默认读 MYSQL_PORT 环境变量"},
                "user":     {"type": "string",  "description": "默认读 MYSQL_USER 环境变量"},
                "password": {"type": "string",  "description": "默认读 MYSQL_PASSWORD 环境变量"},
                "database": {"type": "string",  "description": "默认读 MYSQL_DATABASE 环境变量"},
                "table":    {"type": "string",  "description": "不填则列出所有表；填写表名则返回 DESCRIBE 结果"},
                "sql":      {"type": "string",  "description": "直接执行的 SQL（如 DESCRIBE/SHOW/SELECT FROM INFORMATION_SCHEMA），优先于 table 参数"},
            },
        }

    async def run(self, arguments: dict[str, Any]) -> dict[str, Any]:
        import aiomysql
        cfg = _mysql_defaults(arguments)
        raw_sql = (arguments.get("sql") or "").strip()
        table   = arguments.get("table", "")
        if raw_sql:
            query = raw_sql
        elif table:
            query = f"DESCRIBE `{table}`"
        else:
            query = "SHOW TABLES"
        conn = await aiomysql.connect(
            host=cfg["host"], port=cfg["port"], user=cfg["user"],
            password=cfg["password"], db=cfg["database"],
            charset="utf8mb4",
        )
        try:
            async with conn.cursor(aiomysql.DictCursor) as cur:
                await cur.execute(query)
                rows = await cur.fetchall()
                return {"rows": [dict(r) for r in rows]}
        except Exception as e:
            return {"error": str(e)}
        finally:
            conn.close()


class RedisCmdTool(BaseTool):
    name = "redis_cmd"
    description = "执行 Redis 命令，如 GET key / SET key value / DEL key / KEYS pattern / HGETALL key 等"
    category = "database"
    risk_level = "medium"

    @property
    def input_schema(self) -> dict:
        return {
            "type": "object",
            "properties": {
                "host":     {"type": "string", "default": "127.0.0.1"},
                "port":     {"type": "integer", "default": 6379},
                "password": {"type": "string", "default": ""},
                "db":       {"type": "integer", "default": 0},
                "command":  {"type": "string", "description": "Redis 命令，如 GET / SET / DEL / KEYS / HGETALL"},
                "args":     {"type": "array", "items": {"type": "string"}, "description": "命令参数列表"},
            },
            "required": ["command"],
        }

    async def run(self, arguments: dict[str, Any]) -> dict[str, Any]:
        import redis.asyncio as aioredis
        host = arguments.get("host", "127.0.0.1")
        port = int(arguments.get("port", 6379))
        password = arguments.get("password") or None
        db = int(arguments.get("db", 0))
        command = arguments["command"].upper()
        args = arguments.get("args", [])
        try:
            r = aioredis.Redis(host=host, port=port, password=password, db=db,
                               decode_responses=True)
            result = await r.execute_command(command, *args)
            await r.aclose()
            if isinstance(result, (list, dict)):
                return {"result": result}
            return {"result": str(result) if result is not None else None}
        except Exception as e:
            return {"error": str(e)}


class SQLiteQueryTool(BaseTool):
    name = "sqlite_query"
    description = "查询本地 SQLite 数据库文件，执行 SQL 语句（无需服务器）"
    category = "database"
    risk_level = "medium"

    @property
    def input_schema(self) -> dict:
        return {
            "type": "object",
            "properties": {
                "db_path":  {"type": "string", "description": "SQLite 文件路径"},
                "sql":      {"type": "string"},
                "max_rows": {"type": "integer", "default": 200},
            },
            "required": ["db_path", "sql"],
        }

    async def run(self, arguments: dict[str, Any]) -> dict[str, Any]:
        import sqlite3
        import asyncio
        db_path = arguments["db_path"]
        sql = arguments["sql"]
        max_rows = int(arguments.get("max_rows", 200))

        def _run():
            conn = sqlite3.connect(db_path)
            conn.row_factory = sqlite3.Row
            cur = conn.execute(sql)
            if sql.strip().upper().startswith("SELECT") or sql.strip().upper().startswith("PRAGMA"):
                rows = [dict(r) for r in cur.fetchmany(max_rows)]
                conn.close()
                return {"rows": rows, "count": len(rows)}
            else:
                conn.commit()
                affected = cur.rowcount
                conn.close()
                return {"affected_rows": affected, "ok": True}

        try:
            loop = asyncio.get_event_loop()
            return await loop.run_in_executor(None, _run)
        except Exception as e:
            return {"error": str(e)}


# ─────────────────────────────────────────────────────────────────────────────
# H. Office 文档
# ─────────────────────────────────────────────────────────────────────────────

class ExcelReadTool(BaseTool):
    name = "excel_read"
    description = "读取 Excel (.xlsx/.xls) 文件，返回表格数据"
    category = "document"
    risk_level = "low"

    @property
    def input_schema(self) -> dict:
        return {
            "type": "object",
            "properties": {
                "path":     {"type": "string", "description": "Excel 文件路径"},
                "sheet":    {"type": "string", "description": "Sheet 名称，不填则读第一个"},
                "max_rows": {"type": "integer", "default": 500},
            },
            "required": ["path"],
        }

    async def run(self, arguments: dict[str, Any]) -> dict[str, Any]:
        import asyncio
        path = arguments["path"]
        sheet = arguments.get("sheet")
        max_rows = int(arguments.get("max_rows", 500))

        def _read():
            import openpyxl
            wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
            ws = wb[sheet] if sheet else wb.active
            headers: list = []
            rows: list = []
            for i, row in enumerate(ws.iter_rows(values_only=True)):
                if i == 0:
                    headers = [str(c) if c is not None else f"col{j}" for j, c in enumerate(row)]
                else:
                    rows.append(dict(zip(headers, row)))
                if i >= max_rows:
                    break
            wb.close()
            return {"sheet": ws.title, "headers": headers, "rows": rows, "count": len(rows)}

        try:
            return await asyncio.get_event_loop().run_in_executor(None, _read)
        except Exception as e:
            return {"error": str(e)}


class ExcelWriteTool(BaseTool):
    name = "excel_write"
    description = "将数据写入 Excel 文件（列表的列表或 dict 列表）"
    category = "document"
    risk_level = "medium"

    @property
    def input_schema(self) -> dict:
        return {
            "type": "object",
            "properties": {
                "path":    {"type": "string"},
                "data":    {"type": "array", "description": "二维数组或 dict 列表"},
                "sheet":   {"type": "string", "default": "Sheet1"},
                "headers": {"type": "array", "items": {"type": "string"}, "description": "列标题（dict列表时可省略）"},
            },
            "required": ["path", "data"],
        }

    async def run(self, arguments: dict[str, Any]) -> dict[str, Any]:
        import asyncio
        path = arguments["path"]
        data = arguments["data"]
        sheet_name = arguments.get("sheet", "Sheet1")
        headers = arguments.get("headers")

        def _write():
            import openpyxl
            wb = openpyxl.Workbook()
            ws = wb.active
            ws.title = sheet_name
            if data and isinstance(data[0], dict):
                h = headers or list(data[0].keys())
                ws.append(h)
                for row in data:
                    ws.append([row.get(k) for k in h])
            else:
                if headers:
                    ws.append(headers)
                for row in data:
                    ws.append(list(row))
            wb.save(path)
            return {"ok": True, "path": path, "rows_written": len(data)}

        try:
            return await asyncio.get_event_loop().run_in_executor(None, _write)
        except Exception as e:
            return {"error": str(e)}


class WordReadTool(BaseTool):
    name = "word_read"
    description = "读取 Word (.docx) 文档，提取所有段落文本和表格内容"
    category = "document"
    risk_level = "low"

    @property
    def input_schema(self) -> dict:
        return {
            "type": "object",
            "properties": {
                "path": {"type": "string", "description": "Word 文件路径"},
            },
            "required": ["path"],
        }

    async def run(self, arguments: dict[str, Any]) -> dict[str, Any]:
        import asyncio

        def _read():
            from docx import Document
            doc = Document(arguments["path"])
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            tables = []
            for t in doc.tables:
                table_data = [[cell.text for cell in row.cells] for row in t.rows]
                tables.append(table_data)
            return {
                "paragraphs": paragraphs,
                "full_text": "\n".join(paragraphs),
                "tables": tables,
                "table_count": len(tables),
            }

        try:
            return await asyncio.get_event_loop().run_in_executor(None, _read)
        except Exception as e:
            return {"error": str(e)}


class WordWriteTool(BaseTool):
    name = "word_write"
    description = "创建 Word (.docx) 文档，支持标题和段落"
    category = "document"
    risk_level = "medium"

    @property
    def input_schema(self) -> dict:
        return {
            "type": "object",
            "properties": {
                "path":    {"type": "string"},
                "title":   {"type": "string", "description": "文档标题（可选）"},
                "content": {"type": "string", "description": "正文内容，用 \\n\\n 分隔段落"},
            },
            "required": ["path", "content"],
        }

    async def run(self, arguments: dict[str, Any]) -> dict[str, Any]:
        import asyncio

        def _write():
            from docx import Document
            from docx.shared import Pt
            doc = Document()
            if arguments.get("title"):
                doc.add_heading(arguments["title"], level=1)
            for para in arguments["content"].split("\n\n"):
                para = para.strip()
                if para.startswith("# "):
                    doc.add_heading(para[2:], level=2)
                elif para.startswith("## "):
                    doc.add_heading(para[3:], level=3)
                elif para:
                    doc.add_paragraph(para)
            doc.save(arguments["path"])
            return {"ok": True, "path": arguments["path"]}

        try:
            return await asyncio.get_event_loop().run_in_executor(None, _write)
        except Exception as e:
            return {"error": str(e)}


class PPTReadTool(BaseTool):
    name = "ppt_read"
    description = "读取 PowerPoint (.pptx) 文件，提取每张幻灯片的标题和文本内容"
    category = "document"
    risk_level = "low"

    @property
    def input_schema(self) -> dict:
        return {
            "type": "object",
            "properties": {
                "path": {"type": "string", "description": "PPT 文件路径"},
            },
            "required": ["path"],
        }

    async def run(self, arguments: dict[str, Any]) -> dict[str, Any]:
        import asyncio

        def _read():
            from pptx import Presentation
            prs = Presentation(arguments["path"])
            slides = []
            for i, slide in enumerate(prs.slides, 1):
                texts = []
                title = ""
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        t = shape.text_frame.text.strip()
                        if not title and shape.shape_type == 13:
                            title = t
                        elif t:
                            texts.append(t)
                    if hasattr(shape, "name") and "title" in shape.name.lower() and shape.has_text_frame:
                        title = shape.text_frame.text.strip()
                slides.append({"slide": i, "title": title, "content": "\n".join(texts)})
            return {"slides": slides, "total": len(slides)}

        try:
            return await asyncio.get_event_loop().run_in_executor(None, _read)
        except Exception as e:
            return {"error": str(e)}


class PPTWriteTool(BaseTool):
    name = "ppt_write"
    description = "创建 PowerPoint (.pptx) 文件，每个元素为一张幻灯片 {title, content}"
    category = "document"
    risk_level = "medium"

    @property
    def input_schema(self) -> dict:
        return {
            "type": "object",
            "properties": {
                "path":   {"type": "string"},
                "slides": {
                    "type": "array",
                    "items": {
                        "type": "object",
                        "properties": {
                            "title":   {"type": "string"},
                            "content": {"type": "string"},
                        },
                    },
                    "description": "幻灯片列表，每项含 title 和 content",
                },
            },
            "required": ["path", "slides"],
        }

    async def run(self, arguments: dict[str, Any]) -> dict[str, Any]:
        import asyncio

        def _write():
            from pptx import Presentation
            from pptx.util import Inches, Pt
            prs = Presentation()
            layout = prs.slide_layouts[1]  # title + content
            for slide_data in arguments["slides"]:
                slide = prs.slides.add_slide(layout)
                slide.shapes.title.text = slide_data.get("title", "")
                body = slide.placeholders[1]
                body.text = slide_data.get("content", "")
            prs.save(arguments["path"])
            return {"ok": True, "path": arguments["path"], "slides": len(arguments["slides"])}

        try:
            return await asyncio.get_event_loop().run_in_executor(None, _write)
        except Exception as e:
            return {"error": str(e)}


class MdToHtmlTool(BaseTool):
    name = "md_to_html"
    description = "将 Markdown 内容转换为 HTML（支持表格、代码高亮等 GFM 特性）"
    category = "document"
    risk_level = "low"

    @property
    def input_schema(self) -> dict:
        return {
            "type": "object",
            "properties": {
                "content":     {"type": "string", "description": "Markdown 文本（与 path 二选一）"},
                "path":        {"type": "string", "description": "Markdown 文件路径（与 content 二选一）"},
                "output_path": {"type": "string", "description": "输出 HTML 文件路径（不填则只返回字符串）"},
            },
        }

    async def run(self, arguments: dict[str, Any]) -> dict[str, Any]:
        try:
            import markdown
            content = arguments.get("content", "")
            if not content and arguments.get("path"):
                with open(arguments["path"], encoding="utf-8") as f:
                    content = f.read()
            html = markdown.markdown(
                content,
                extensions=["tables", "fenced_code", "codehilite", "toc"],
            )
            if arguments.get("output_path"):
                with open(arguments["output_path"], "w", encoding="utf-8") as f:
                    f.write(f"<!DOCTYPE html><html><body>{html}</body></html>")
                return {"ok": True, "output_path": arguments["output_path"]}
            return {"html": html}
        except ImportError:
            import re
            content = arguments.get("content", "")
            html = re.sub(r"\*\*(.*?)\*\*", r"<b>\1</b>", content)
            html = re.sub(r"\n", "<br>", html)
            return {"html": html, "warning": "markdown 包未安装，使用简易转换"}
        except Exception as e:
            return {"error": str(e)}


# Milvus (optional — only registers if pymilvus is available)
class MilvusSearchTool(BaseTool):
    name = "milvus_search"
    description = "在 Milvus 向量数据库中执行相似度搜索（需要运行中的 Milvus 服务）"
    category = "database"
    risk_level = "low"

    @property
    def input_schema(self) -> dict:
        return {
            "type": "object",
            "properties": {
                "host":       {"type": "string", "default": "localhost"},
                "port":       {"type": "integer", "default": 19530},
                "collection": {"type": "string"},
                "vector":     {"type": "array", "items": {"type": "number"}, "description": "查询向量"},
                "field":      {"type": "string", "description": "向量字段名", "default": "embedding"},
                "top_k":      {"type": "integer", "default": 5},
                "output_fields": {"type": "array", "items": {"type": "string"}},
            },
            "required": ["collection", "vector"],
        }

    async def run(self, arguments: dict[str, Any]) -> dict[str, Any]:
        import asyncio

        def _search():
            from pymilvus import connections, Collection
            connections.connect(
                host=arguments.get("host", "localhost"),
                port=arguments.get("port", 19530),
            )
            col = Collection(arguments["collection"])
            col.load()
            results = col.search(
                data=[arguments["vector"]],
                anns_field=arguments.get("field", "embedding"),
                param={"metric_type": "L2", "params": {"nprobe": 10}},
                limit=arguments.get("top_k", 5),
                output_fields=arguments.get("output_fields"),
            )
            hits = [{"id": h.id, "distance": h.distance, "entity": h.entity.to_dict()}
                    for h in results[0]]
            return {"hits": hits}

        try:
            return await asyncio.get_event_loop().run_in_executor(None, _search)
        except Exception as e:
            return {"error": str(e)}


# Export list for registry
ALL_DATA_TOOLS = [
    MySQLQueryTool, MySQLSchemaTool,
    RedisCmdTool,
    SQLiteQueryTool,
    ExcelReadTool, ExcelWriteTool,
    WordReadTool, WordWriteTool,
    PPTReadTool, PPTWriteTool,
    MdToHtmlTool,
    MilvusSearchTool,
]
